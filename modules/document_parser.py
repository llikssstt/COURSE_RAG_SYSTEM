from __future__ import annotations

from pathlib import Path
from typing import Any


def _clean_text(text: str) -> str:
    return "\n".join(line.strip() for line in (text or "").splitlines() if line.strip())


def parse_pdf(path: Path, file_name: str) -> list[dict[str, Any]]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("缺少 pypdf，请先安装 requirements.txt") from exc

    reader = PdfReader(str(path))
    items = []
    for idx, page in enumerate(reader.pages, start=1):
        text = _clean_text(page.extract_text() or "")
        if text:
            items.append(
                {
                    "text": text,
                    "source_file": file_name,
                    "page_number": idx,
                    "slide_number": None,
                    "section_title": f"第 {idx} 页",
                }
            )
    return items


def parse_docx(path: Path, file_name: str) -> list[dict[str, Any]]:
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("缺少 python-docx，请先安装 requirements.txt") from exc

    doc = Document(str(path))
    items = []
    current_title = ""
    para_idx = 0
    for para in doc.paragraphs:
        text = _clean_text(para.text)
        if not text:
            continue
        para_idx += 1
        style_name = para.style.name if para.style else ""
        if "Heading" in style_name or "标题" in style_name:
            current_title = text
        items.append(
            {
                "text": text,
                "source_file": file_name,
                "page_number": None,
                "slide_number": None,
                "section_title": current_title or f"段落 {para_idx}",
            }
        )
    return items


def parse_pptx(path: Path, file_name: str) -> list[dict[str, Any]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError("缺少 python-pptx，请先安装 requirements.txt") from exc

    prs = Presentation(str(path))
    items = []
    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = _clean_text(shape.text)
                if text:
                    texts.append(text)
        if not texts:
            continue
        title = texts[0][:80]
        items.append(
            {
                "text": "\n".join(texts),
                "source_file": file_name,
                "page_number": None,
                "slide_number": slide_idx,
                "section_title": title or f"幻灯片 {slide_idx}",
            }
        )
    return items


def parse_document(path: str | Path) -> list[dict[str, Any]]:
    file_path = Path(path)
    suffix = file_path.suffix.lower()
    file_name = file_path.name
    if suffix == ".pdf":
        return parse_pdf(file_path, file_name)
    if suffix == ".docx":
        return parse_docx(file_path, file_name)
    if suffix == ".pptx":
        return parse_pptx(file_path, file_name)
    raise ValueError(f"不支持的文件类型：{suffix}")
